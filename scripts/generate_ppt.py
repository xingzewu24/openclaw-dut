#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
大连理工大学 PPT 生成工具
功能：基于模板生成 PPT，支持从 Markdown 自动生成多页幻灯片
依赖：python-pptx (pip install python-pptx)
模板目录：~/.openclaw/workspace/skills/dlut-campus/templates/
"""

import os
import sys
import re
import platform
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("❌ 缺少 python-pptx 库，请执行: pip install python-pptx")
    sys.exit(1)

SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATES_DIR = os.path.join(SKILL_DIR, "templates")
DEFAULT_TEMPLATE = os.path.join(TEMPLATES_DIR, "0.大连理工大学通用PPT模板.pptx")
DUT_BLUE = RGBColor(0, 51, 102)
DARK_TEXT = RGBColor(45, 45, 45)
MUTED_TEXT = RGBColor(110, 110, 110)

# 跨平台字体：macOS 用 PingFang SC，Windows 用微软雅黑
_DEFAULT_FONT = "Microsoft YaHei" if platform.system() == "Windows" else "PingFang SC"


def list_templates():
    """列出所有可用模板"""
    templates = []
    if not os.path.isdir(TEMPLATES_DIR):
        return templates
    for f in sorted(os.listdir(TEMPLATES_DIR)):
        if f.lower().endswith((".pptx", ".ppt")) and not f.startswith("~$"):
            fpath = os.path.join(TEMPLATES_DIR, f)
            size_kb = os.path.getsize(fpath) / 1024
            templates.append({
                "name": f,
                "path": fpath,
                "size": f"{size_kb:.1f} KB",
            })
    return templates


def _get_template_path(template_name=None):
    """解析模板路径"""
    if template_name is None:
        if os.path.exists(DEFAULT_TEMPLATE):
            return DEFAULT_TEMPLATE
        return None  # 无模板，创建空白 PPT

    # 精确路径
    if os.path.exists(template_name):
        return template_name

    # 在模板目录中搜索
    if os.path.isdir(TEMPLATES_DIR):
        for f in os.listdir(TEMPLATES_DIR):
            if template_name.lower() in f.lower() and f.lower().endswith(".pptx"):
                return os.path.join(TEMPLATES_DIR, f)

    return None


def _find_layout_by_names(prs, candidate_names):
    """按给定名称顺序优先匹配版式"""
    for candidate in candidate_names:
        for layout in prs.slide_layouts:
            name = getattr(layout, "name", "") or ""
            if name == candidate:
                return layout
    return None


def _get_layout(prs, layout_name):
    """获取幻灯片版式，优先按名称匹配，兼容通用模板与大工模板"""
    name_map = {
        "title": ["封面", "Title Slide", "标题幻灯片"],
        "content": ["标题和内容", "Title and Content", "仅标题"],
        "section": ["章节过渡页", "章节过渡页-2", "Section Header"],
        "blank": ["空白", "空白（lowpoly）", "Blank"],
    }
    layout = _find_layout_by_names(prs, name_map.get(layout_name, []))
    if layout is not None:
        return layout

    # 名称匹配失败时回退到旧索引策略
    layout_map = {
        "title": 0,
        "content": 1,
        "section": 2,
        "blank": 6,
    }
    idx = layout_map.get(layout_name, 1)
    layouts = prs.slide_layouts
    if idx < len(layouts):
        return layouts[idx]
    return layouts[0] if layouts else None


def _clear_all_slides(prs):
    """清空模板中的示例页，仅保留母版与版式。"""
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rId = slide_id.rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(slide_id)


def _add_text_to_placeholder(placeholder, text, font_size=20):
    """向占位符添加文本，并优化字号、行距、留白。"""
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    try:
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
    except Exception:
        pass

    lines = [line.strip() for line in text.split("\n") if line.strip()]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(6)
        p.space_before = Pt(0)
        p.line_spacing = 1.15
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(font_size)


def _set_run_style(run, size=None, bold=None, color=None, name=None):
    """统一设置 run 样式。"""
    if name is None:
        name = _DEFAULT_FONT
    """统一设置 run 样式。"""
    font = run.font
    if size is not None:
        font.size = Pt(size)
    if bold is not None:
        font.bold = bold
    if color is not None:
        font.color.rgb = color
    if name:
        font.name = name


def _get_title_placeholder(slide):
    """获取标题占位符，兼容自定义模板。"""
    if slide.shapes.title is not None:
        return slide.shapes.title
    for shape in slide.placeholders:
        name = getattr(shape, 'name', '') or ''
        if '标题' in name or 'Title' in name:
            return shape
    return None


def _get_body_placeholders(slide):
    """获取正文类占位符列表，按面积从大到小排序。"""
    title_ph = _get_title_placeholder(slide)
    body = []
    for shape in slide.placeholders:
        if shape == title_ph:
            continue
        if hasattr(shape, 'text_frame'):
            body.append(shape)
    body.sort(key=lambda sh: sh.width * sh.height, reverse=True)
    return body


def _get_surrogate_title_placeholder(slide, body_placeholders):
    """当模板没有标准标题占位符时，尝试从正文占位符中找一个最像标题位的。"""
    if not body_placeholders:
        return None
    candidates = sorted(body_placeholders, key=lambda sh: (sh.top, sh.height, -(sh.width * sh.height)))
    for sh in candidates:
        if sh.top < Inches(1.2) and sh.height < Inches(1.2):
            return sh
    return None


def _style_title_shape(shape, size=24):
    """优化标题文字样式。"""
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        for r in p.runs:
            _set_run_style(r, size=size, bold=True, color=DUT_BLUE)


def _style_cover_slide(slide):
    """优化封面页样式。"""
    for sh in slide.shapes:
        if not hasattr(sh, "text_frame") or not sh.text.strip():
            continue
        txt = sh.text.strip()
        tf = sh.text_frame
        for p in tf.paragraphs:
            if not p.runs:
                continue
            for r in p.runs:
                if txt == "大连理工大学":
                    _set_run_style(r, size=17, bold=False, color=MUTED_TEXT)
                else:
                    _set_run_style(r, size=28, bold=True, color=DUT_BLUE)
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(0)


def _style_section_slide(slide):
    """优化章节过渡页样式。"""
    for sh in slide.shapes:
        if not hasattr(sh, "text_frame") or not sh.text.strip():
            continue
        tf = sh.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                _set_run_style(r, size=28, bold=True, color=DUT_BLUE)
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(0)


def _style_body_shape(shape, base_size=19):
    """优化正文字号、层级和强调色。"""
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
    except Exception:
        pass

    original_lines = []
    for p in tf.paragraphs:
        if p.text.strip():
            original_lines.append(p.text)
    if not original_lines and shape.text.strip():
        original_lines = [x for x in shape.text.split("\n") if x.strip()]

    tf.clear()
    for idx, raw in enumerate(original_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(8)
        p.line_spacing = 1.18
        p.alignment = PP_ALIGN.LEFT
        p.level = 0

        text = raw.strip()
        if not text:
            continue

        bullet = ""
        if text.startswith("•"):
            bullet = "• "
            text = text[1:].strip()

        if bullet:
            r0 = p.add_run()
            r0.text = bullet
            _set_run_style(r0, size=base_size, bold=False, color=DARK_TEXT)

        split_idx = -1
        split_char = None
        for ch in ["：", ":", "，", ",", " - ", "（"]:
            pos = text.find(ch)
            if pos != -1 and pos <= 14:
                split_idx = pos
                split_char = ch
                break

        if split_idx != -1:
            if split_char == "（":
                lead = text[:split_idx]
                tail = text[split_idx:]
            else:
                lead = text[:split_idx + len(split_char)]
                tail = text[split_idx + len(split_char):].strip()
            r1 = p.add_run()
            r1.text = lead
            _set_run_style(r1, size=base_size, bold=True, color=DUT_BLUE)
            if tail:
                r2 = p.add_run()
                r2.text = (" " if not tail.startswith("（") else "") + tail
                _set_run_style(r2, size=base_size, bold=False, color=DARK_TEXT)
        else:
            r = p.add_run()
            r.text = text
            _set_run_style(r, size=base_size, bold=False, color=DARK_TEXT)


def _apply_visual_polish(prs):
    """统一修正文档视觉效果，针对大工模板做样式优化。"""
    for slide in prs.slides:
        layout = getattr(slide.slide_layout, "name", "") or ""
        if "封面" in layout:
            _style_cover_slide(slide)
            continue
        if "章节过渡页" in layout:
            _style_section_slide(slide)
            continue

        nonempty = [sh for sh in slide.shapes if hasattr(sh, "text_frame") and sh.text.strip()]
        if not nonempty:
            continue

        title_shape = None
        body_shape = None
        candidates = sorted(nonempty, key=lambda sh: (sh.top, sh.height))
        for sh in candidates:
            if sh.top < Inches(1.2) and sh.height < Inches(1.2):
                title_shape = sh
                break

        area_sorted = sorted(nonempty, key=lambda sh: sh.width * sh.height, reverse=True)
        if area_sorted:
            body_shape = area_sorted[0]
            if body_shape == title_shape and len(area_sorted) > 1:
                body_shape = area_sorted[1]

        if title_shape is not None:
            _style_title_shape(title_shape)
        if body_shape is not None:
            _style_body_shape(body_shape)


def generate_ppt(title, slides_content, template_path=None, output_path="output.pptx", polish=True):
    """
    生成 PPT

    参数:
        title: PPT 标题（用于首页）
        slides_content: 幻灯片内容列表，每项为 dict:
            {"title": "标题", "content": "正文", "layout": "title/content/section/blank"}
        template_path: 模板文件路径（None 则使用默认模板）
        output_path: 输出路径
    """
    tpl = _get_template_path(template_path)

    if tpl and os.path.exists(tpl):
        prs = Presentation(tpl)
        _clear_all_slides(prs)
    else:
        prs = Presentation()
        if tpl is not None:
            print(f"⚠️  模板未找到: {tpl}，使用空白模板")

    # 添加标题页
    title_layout = _get_layout(prs, "title")
    if title_layout:
        slide = prs.slides.add_slide(title_layout)
        title_ph = _get_title_placeholder(slide)
        if title_ph is not None:
            title_ph.text = title
        body_phs = _get_body_placeholders(slide)
        if body_phs:
            _add_text_to_placeholder(body_phs[0], "大连理工大学", font_size=20)

    # 添加内容页
    for item in slides_content:
        layout_name = item.get("layout", "content")
        layout = _get_layout(prs, layout_name)
        if not layout:
            continue

        slide = prs.slides.add_slide(layout)
        placeholders = slide.placeholders

        slide_title = item.get("title", "")
        slide_content = item.get("content", "")

        if layout_name == "blank":
            # 空白页：如果有内容，添加文本框
            if slide_content:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
                tf = txBox.text_frame
                tf.word_wrap = True
                for i, line in enumerate(slide_content.split("\n")):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line.strip()
        else:
            # 有占位符的版式
            body_phs = _get_body_placeholders(slide)
            title_ph = _get_title_placeholder(slide)
            surrogate_title_ph = None

            if title_ph is None:
                surrogate_title_ph = _get_surrogate_title_placeholder(slide, body_phs)
                if surrogate_title_ph is not None:
                    body_phs = [ph for ph in body_phs if ph != surrogate_title_ph]

            if title_ph is not None and slide_title:
                title_ph.text = slide_title
            elif surrogate_title_ph is not None and slide_title:
                _add_text_to_placeholder(surrogate_title_ph, slide_title, font_size=24)
                if surrogate_title_ph.text_frame.paragraphs and surrogate_title_ph.text_frame.paragraphs[0].runs:
                    surrogate_title_ph.text_frame.paragraphs[0].runs[0].font.bold = True
            elif slide_title:
                title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.45), Inches(8.2), Inches(0.7))
                tf = title_box.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_title
                if p.runs:
                    p.runs[0].font.size = Pt(24)
                    p.runs[0].font.bold = True

            if body_phs and slide_content:
                _add_text_to_placeholder(body_phs[0], slide_content, font_size=20)
            elif slide_content:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.8))
                tf = txBox.text_frame
                tf.word_wrap = True
                for i, line in enumerate(slide_content.split("\n")):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line.strip()

    if polish:
        _apply_visual_polish(prs)

    # 确保输出目录存在
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    return {"success": True, "path": os.path.abspath(output_path), "slides": len(slides_content) + 1}


def generate_from_markdown(title, markdown_text, template_path=None, output_path="output.pptx", polish=True):
    """
    从 Markdown 文本生成 PPT
    按 ## 标题切分为各页，支持 # 作为节标题

    参数:
        title: PPT 标题
        markdown_text: Markdown 文本内容
        template_path: 模板路径
        output_path: 输出路径
    """
    slides_content = []

    # 如果是文件路径，读取文件
    if os.path.isfile(markdown_text):
        with open(markdown_text, "r", encoding="utf-8") as f:
            markdown_text = f.read()

    # 按标题切分
    # 先按 ## 切分（二级标题作为幻灯片标题）
    sections = re.split(r'^(#{1,2})\s+(.+)$', markdown_text, flags=re.MULTILINE)

    if len(sections) <= 1:
        # 没有标题标记，整个内容作为一页
        slides_content.append({
            "title": title,
            "content": markdown_text.strip(),
            "layout": "content",
        })
    else:
        # 开头如果有内容（在第一个标题之前）
        preamble = sections[0].strip()
        if preamble:
            slides_content.append({
                "title": "概述",
                "content": _clean_markdown(preamble),
                "layout": "content",
            })

        # 解析标题和内容
        i = 1
        while i < len(sections) - 2:
            level = sections[i]      # # 或 ##
            heading = sections[i+1]  # 标题文本
            content = sections[i+2].strip() if i+2 < len(sections) else ""
            i += 3

            if level == "#":
                layout = "section"
            else:
                layout = "content"

            slides_content.append({
                "title": heading.strip(),
                "content": _clean_markdown(content),
                "layout": layout,
            })

    return generate_ppt(title, slides_content, template_path, output_path, polish=polish)


def _clean_markdown(text):
    """清理 Markdown 标记，转为纯文本"""
    # 移除图片
    text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
    # 链接保留文本
    text = re.sub(r'\[(.+?)\]\(.*?\)', r'\1', text)
    # 移除加粗/斜体标记
    text = re.sub(r'\*{1,2}(.+?)\*{1,2}', r'\1', text)
    # 列表项保留
    text = re.sub(r'^[-*+]\s+', '• ', text, flags=re.MULTILINE)
    # 移除代码块标记
    text = re.sub(r'```\w*\n?', '', text)
    # 移除行内代码标记
    text = re.sub(r'`(.+?)`', r'\1', text)
    # 清理多余空行
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def main():
    parser = argparse.ArgumentParser(description="大连理工大学 PPT 生成工具")
    parser.add_argument("--title", "-t", help="PPT 标题")
    parser.add_argument("--markdown", "-m", help="Markdown 内容或文件路径")
    parser.add_argument("--template", help="模板名称或路径 (默认: 大工通用模板)")
    parser.add_argument("--output", "-o", default="output.pptx", help="输出文件路径")
    parser.add_argument("--no-polish", action="store_true", help="关闭默认的文字样式优化")
    parser.add_argument("--list-templates", action="store_true", help="列出所有可用模板")
    args = parser.parse_args()

    if args.list_templates:
        templates = list_templates()
        if not templates:
            print(f"📂 模板目录为空: {TEMPLATES_DIR}")
        else:
            print(f"\n📂 可用模板 ({TEMPLATES_DIR})")
            print("=" * 50)
            for t in templates:
                print(f"  📄 {t['name']}  ({t['size']})")
        print()
        return

    if not args.title:
        print("❌ 错误: 请提供 --title 参数")
        sys.exit(1)

    if args.markdown:
        result = generate_from_markdown(
            title=args.title,
            markdown_text=args.markdown,
            template_path=args.template,
            output_path=args.output,
            polish=not args.no_polish,
        )
    else:
        # 无 markdown 时创建只有标题页的 PPT
        result = generate_ppt(
            title=args.title,
            slides_content=[],
            template_path=args.template,
            output_path=args.output,
            polish=not args.no_polish,
        )

    if result["success"]:
        print(f"✅ PPT 生成成功!")
        print(f"   📄 文件: {result['path']}")
        print(f"   📊 页数: {result['slides']}")
    else:
        print("❌ PPT 生成失败")
        sys.exit(1)


if __name__ == "__main__":
    main()
