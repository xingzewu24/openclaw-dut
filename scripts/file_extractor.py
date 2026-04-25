#!/usr/bin/env python3
"""课件内容提取器 - PPT/PDF/DOCX → 纯文本/Markdown"""

import os
import sys
from pathlib import Path

def extract_pptx(file_path):
    """从 PPTX 提取文本内容"""
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)
        if texts:
            slides.append(f"## Slide {i}\n\n" + "\n".join(f"- {t}" for t in texts))
    return "\n\n".join(slides)

def extract_pdf(file_path):
    """从 PDF 提取文本内容"""
    import pdfplumber
    texts = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                texts.append(f"## Page {i}\n\n{text.strip()}")
    return "\n\n".join(texts)

def extract_docx(file_path):
    """从 DOCX 提取文本内容"""
    try:
        from docx import Document
        doc = Document(file_path)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())
        return "\n\n".join(texts)
    except ImportError:
        return "[需要安装 python-docx: pip3 install python-docx]"

def extract_file(file_path):
    """根据文件类型自动提取文本"""
    ext = Path(file_path).suffix.lower()
    extractors = {
        ".pptx": extract_pptx,
        ".ppt": extract_pptx,
        ".pdf": extract_pdf,
        ".docx": extract_docx,
        ".txt": lambda f: open(f).read(),
        ".md": lambda f: open(f).read(),
    }
    extractor = extractors.get(ext)
    if not extractor:
        return f"[不支持的文件类型: {ext}]"
    try:
        return extractor(file_path)
    except Exception as e:
        return f"[提取失败: {e}]"

def extract_to_markdown(file_path, output_path=None):
    """提取文件内容并保存为 Markdown"""
    content = extract_file(file_path)
    fname = Path(file_path).stem
    md = f"# {fname}\n\n{content}"
    if output_path:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(md)
    return md

def batch_extract(directory, output_dir=None, extensions=None):
    """批量提取目录下的所有课件"""
    if extensions is None:
        extensions = {".pptx", ".pdf", ".docx"}
    results = []
    for root, dirs, files in os.walk(directory):
        for fname in sorted(files):
            ext = Path(fname).suffix.lower()
            if ext not in extensions:
                continue
            fpath = os.path.join(root, fname)
            if output_dir:
                rel = os.path.relpath(fpath, directory)
                out_path = os.path.join(output_dir, Path(rel).with_suffix(".md"))
                md = extract_to_markdown(fpath, out_path)
                print(f"✅ {fname} → {out_path}")
            else:
                md = extract_file(fpath)
                print(f"✅ {fname} ({len(md)} chars)")
            results.append({"file": fname, "content": md})
    return results

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: file_extractor.py <file_or_dir> [output_dir]")
        sys.exit(1)
    
    target = sys.argv[1]
    output = sys.argv[2] if len(sys.argv) > 2 else None
    
    if os.path.isdir(target):
        batch_extract(target, output)
    else:
        content = extract_file(target)
        print(content[:2000])
        if len(content) > 2000:
            print(f"\n... (共 {len(content)} 字符)")
